import os
import logging
from flask import Flask, request, jsonify, render_template
from flask_cors import CORS
import fitz  # PyMuPDF
from transformers import pipeline
from collections import defaultdict
from sumy.parsers.plaintext import PlaintextParser
from sumy.nlp.tokenizers import Tokenizer
from sumy.summarizers.lex_rank import LexRankSummarizer
import traceback
import re
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from io import BytesIO
from flask import send_file

# Load environment variables from .env file
load_dotenv()


logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

app = Flask(__name__)
CORS(app)

# Load the summarization model
summarizer = pipeline("summarization", model="facebook/bart-large-cnn")

UPLOAD_FOLDER = 'Uploads'
if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)

def extract_sections(text):
    """
    Splits the document into key sections based on common headings,
    returns sections in fixed order (Abstract, Introduction, Proposed Methodology, Results, Conclusion)
    and skips any sections not present.
    """
    sections_order = [
        "Abstract",
        "Introduction",
        "Proposed Methodology",
        "Results",
        "Conclusion"
    ]
    sections = {title: "" for title in sections_order}

    lower_text = text.lower()
    section_indices = []

    # Find all headings positions (case-insensitive)
    for title in sections_order:
        idx = lower_text.find(title.lower())
        if idx != -1:
            section_indices.append((idx, title))
    # Sort found headings by their position in text
    section_indices.sort(key=lambda x: x[0])

    # Map each heading to start and end position of its content
    section_pos = {}
    for i, (start_idx, title) in enumerate(section_indices):
        # Find end of heading line (skip heading line)
        heading_end = text.find('\n', start_idx)
        if heading_end == -1:
            heading_end = start_idx
        else:
            heading_end += 1  # Start after heading line

        # End of section is start of next heading or end of text
        end_idx = section_indices[i + 1][0] if i + 1 < len(section_indices) else len(text)
        section_pos[title] = (heading_end, end_idx)

    # Extract section text only for found sections
    found_sections = {}
    for title in sections_order:
        if title in section_pos:
            start, end = section_pos[title]
            content = text[start:end].strip()
            if content:
                found_sections[title] = content

    return found_sections  # Return only sections actually found, in the fixed order


def summarize_section(section_text, sentences_count=5):
    if not section_text.strip():
        return "Section not found in the document."
    try:
        parser = PlaintextParser.from_string(section_text, Tokenizer("english"))
        summarizer = LexRankSummarizer()
        summary_sentences = summarizer(parser.document, sentences_count)
        summary = ' '.join(str(sentence) for sentence in summary_sentences)
        return summary if summary else "Summary not available."
    except Exception as e:
        logger.error(f"Summarization error: {e}")
        logger.error(traceback.format_exc())
        return f"Error during summarization: {e}"
    
def summarize_section(section_text):
    """Basic summarization using simple length-based truncation as placeholder."""
    if not section_text.strip():
        return "Section not found in the document."
    lines = section_text.strip().splitlines()
    summary = ' '.join(lines[:5])  # Return first 5 lines as summary placeholder
    return summary if summary else "Summary not available."

def safe_summarize(content):
    """Summarize content with dynamic length based on input size."""
    if not content.strip():
        return ""
    words = len(content.split())
    if words < 30:
        max_len, min_len = 50, 20
    elif words < 80:
        max_len, min_len = 100, 40
    elif words < 200:
        max_len, min_len = 150, 60
    else:
        max_len, min_len = 180, 80
    try:
        summary = summarizer(content, max_length=max_len, min_length=min_len, do_sample=False)[0]['summary_text']
        return summary
    except:
        return ""

def generate_keypoints(section_summaries):
    """Generate keypoints from summarized sections."""
    keypoints = {}
    for heading, summary in section_summaries.items():
        if summary and summary != "Section not found in the document." and summary != "Summary not available.":
            keypoint_summary = safe_summarize(summary)
            if keypoint_summary:
                points = re.split(r'(?<=[.!?]) +', keypoint_summary)
                keypoints[heading] = [point.strip() for point in points if point.strip()]
            else:
                keypoints[heading] = []
        else:
            keypoints[heading] = []
    return keypoints


@app.route('/')
def index():
    try:
        # Pass Firebase config to the template
        firebase_config = {
            'apiKey': os.getenv('FIREBASE_API_KEY'),
            'authDomain': os.getenv('FIREBASE_AUTH_DOMAIN'),
            'projectId': os.getenv('FIREBASE_PROJECT_ID'),
            'storageBucket': os.getenv('FIREBASE_STORAGE_BUCKET'),
            'messagingSenderId': os.getenv('FIREBASE_MESSAGING_SENDER_ID'),
            'appId': os.getenv('FIREBASE_APP_ID'),
            'measurementId': os.getenv('FIREBASE_MEASUREMENT_ID')
        }
        return render_template('index.html', firebase_config=firebase_config)
    except Exception as e:
        return jsonify({'error': 'Template not found'}), 404

@app.route('/customize')
def customize():

    try:
        return render_template('customize.html')
    except Exception as e:
        return jsonify({'error': 'Template not found'}), 404

@app.route('/upload')
def upload():
    try:
        return render_template('upload.html')
    except Exception as e:
        return jsonify({'error': 'Template not found'}), 404

def hex_to_rgb(hex_color):
    """Convert hex color like '#ffcc00' to RGBColor"""
    hex_color = hex_color.lstrip("#")
    return RGBColor(int(hex_color[0:2], 16),
                    int(hex_color[2:4], 16),
                    int(hex_color[4:6], 16))

@app.route('/generate_ppt', methods=['POST'])
def generate_ppt():
    slides_data = request.json['slides']
    prs = Presentation()

    for slide_info in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title_shape = slide.shapes.title
        content_shape = slide.placeholders[1]

        # Set title and bullet points
        title_shape.text = slide_info['title']
        content_shape.text = '\n'.join(slide_info['points'])

        styles = slide_info.get('styles', {})
        font_name = styles.get('font')
        font_color = styles.get('color')
        bg_color = styles.get('backgroundColor')

        # Apply background color
        if bg_color:
            try:
                rgb = hex_to_rgb(bg_color)
                slide.background.fill.solid()
                slide.background.fill.fore_color.rgb = rgb
            except Exception as e:
                print(f"Invalid background color: {bg_color}, error: {e}")

        # Apply font name and color to all text in this slide
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(18)
                        if font_name:
                            run.font.name = font_name
                        if font_color:
                            try:
                                run.font.color.rgb = hex_to_rgb(font_color)
                            except Exception as e:
                                print(f"Invalid font color: {font_color}, error: {e}")

    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return send_file(output, download_name="custom_slides.pptx", as_attachment=True)


@app.route('/process', methods=['POST', 'OPTIONS'])
def process_document():
    if request.method == 'OPTIONS':
        response = app.make_default_options_response()
        response.headers['Access-Control-Allow-Origin'] = '*'
        response.headers['Access-Control-Allow-Methods'] = 'POST, OPTIONS'
        response.headers['Access-Control-Allow-Headers'] = 'Content-Type'
        return response

    if 'file' not in request.files or request.files['file'].filename == '':
        return jsonify({'error': 'No file uploaded'}), 400

    file = request.files['file']
    filename = file.filename
    file_path = os.path.join(UPLOAD_FOLDER, filename)

    try:
        file.save(file_path)

        if filename.lower().endswith('.pdf'):
            try:
                doc = fitz.open(file_path)
                text = "".join(page.get_text() for page in doc)
                doc.close()
            except Exception as e:
                return jsonify({'filename': filename, 'error': str(e)}), 500

        elif filename.lower().endswith('.txt'):
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    text = f.read()
            except Exception as e:
                return jsonify({'filename': filename, 'error': str(e)}), 500

        else:
            return jsonify({'filename': filename, 'error': 'Unsupported file format'}), 400

        # Clean up
        try:
            os.remove(file_path)
        except Exception as e:
            logger.warning(f"Could not remove file: {str(e)}")

        # Section-wise extraction
        section_texts = extract_sections(text)
        section_summaries = {sec: summarize_section(content) for sec, content in section_texts.items()}
        keypoints = generate_keypoints(section_summaries)

        result = {
            'filename': filename,
            'text': text,
            'summary': section_summaries,
            'keypoints': keypoints
        }
        
        return jsonify(result)

    except Exception as e:
        return jsonify({'filename': filename, 'error': str(e)}), 500

if __name__ == '__main__':
    app.run(debug=True, port=8080)
